"""Indexing of DOCX and XLSX documents.

These docs are ZIP-based binaries; the indexer must extract their text so they
become searchable via SQLite FTS5 like code/text/PDF. Fixtures are generated
synthetically (no committed binaries) so the tests are reproducible.
"""
from __future__ import annotations

from pathlib import Path

from fastapi.testclient import TestClient

from services.api.app.config import settings
from services.api.app.main import app

client = TestClient(app)

# Unique markers so a search query hits exactly one document.
DOCX_MARKER = 'NEBULA_QUARTERLY_REPORT_MARKER'
XLSX_MARKER = 'ZIRCON_REVENUE_BREAKDOWN_MARKER'
PPTX_MARKER = 'HALCYON_ROADMAP_SLIDE_MARKER'


def _seed_office_fixtures() -> None:
    """Write a tiny .docx, .xlsx and .pptx into the shared demo project dir."""
    from docx import Document
    import openpyxl
    from pptx import Presentation

    demo = settings.workspace_root / 'demo'
    demo.mkdir(parents=True, exist_ok=True)

    doc = Document()
    doc.add_paragraph(f'This document contains the {DOCX_MARKER} for Q3.')
    doc.add_paragraph('Ordinary filler paragraph with no marker.')
    doc.save(demo / 'report.docx')

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Revenue'
    ws.append(['Region', 'Amount'])
    ws.append(['North', XLSX_MARKER])
    ws.append(['South', '200'])
    wb.save(demo / 'finance.xlsx')

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = f'Roadmap: {PPTX_MARKER}'
    deck.save(demo / 'roadmap.pptx')


def _create_demo_project() -> dict:
    projects = client.get('/api/projects').json()
    for p in projects:
        if p['name'] == 'Demo':
            return p
    response = client.post('/api/projects', json={
        'name': 'Demo', 'path': str(settings.workspace_root / 'demo')
    })
    assert response.status_code == 200, response.text
    return response.json()


def test_docx_and_xlsx_are_indexed_and_searchable() -> None:
    _seed_office_fixtures()
    project = _create_demo_project()
    result = client.post(f"/api/projects/{project['id']}/index")
    assert result.status_code == 200, result.text
    # Both office docs must have been picked up (not silently skipped as binary).
    indexed_files = client.get(f"/api/projects/{project['id']}/files").json()
    rel_paths = {f['relative_path'] for f in indexed_files}
    assert 'report.docx' in rel_paths, f'docx not indexed; got {sorted(rel_paths)}'
    assert 'finance.xlsx' in rel_paths, f'xlsx not indexed; got {sorted(rel_paths)}'

    # The unique markers must be retrievable via FTS5 search.
    docx_hit = client.post('/api/search', json={
        'project_id': project['id'], 'query': DOCX_MARKER
    }).json()
    assert any(h['relative_path'] == 'report.docx' for h in docx_hit['results']), docx_hit

    xlsx_hit = client.post('/api/search', json={
        'project_id': project['id'], 'query': XLSX_MARKER
    }).json()
    assert any(h['relative_path'] == 'finance.xlsx' for h in xlsx_hit['results']), xlsx_hit


def test_pptx_is_indexed_searchable_and_parser_tracked() -> None:
    _seed_office_fixtures()
    project = _create_demo_project()
    result = client.post(f"/api/projects/{project['id']}/index")
    assert result.status_code == 200, result.text

    indexed_files = client.get(f"/api/projects/{project['id']}/files").json()
    rel_paths = {f['relative_path'] for f in indexed_files}
    assert 'roadmap.pptx' in rel_paths, f'pptx not indexed; got {sorted(rel_paths)}'

    # parser/parse_status are internal-only columns (not in the /files API
    # response, same as v2) so a future audit can tell which extractor
    # handled a file without re-reading it. Verify directly against the DB.
    from services.api.app.database import connect
    with connect() as conn:
        row = conn.execute(
            "SELECT parser, parse_status FROM files WHERE project_id=? AND relative_path='roadmap.pptx'",
            (project['id'],)
        ).fetchone()
    assert row is not None
    assert row['parser'] == 'pptx'
    assert row['parse_status'] == 'indexed'

    pptx_hit = client.post('/api/search', json={
        'project_id': project['id'], 'query': PPTX_MARKER
    }).json()
    assert any(h['relative_path'] == 'roadmap.pptx' for h in pptx_hit['results']), pptx_hit


def test_corrupt_office_file_does_not_crash_indexing() -> None:
    # A file with an office extension but garbage content must be skipped, not
    # raise and abort the whole index run.
    demo = settings.workspace_root / 'demo'
    (demo / 'broken.docx').write_bytes(b'not actually a docx zip')
    project = _create_demo_project()
    result = client.post(f"/api/projects/{project['id']}/index")
    assert result.status_code == 200, result.text
    # broken.docx appears in errors list but indexing still completed.
    body = result.json()
    assert any('broken.docx' in e for e in body.get('errors', [])), body
