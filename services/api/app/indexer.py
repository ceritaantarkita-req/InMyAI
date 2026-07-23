from __future__ import annotations

import hashlib
import json
from pathlib import Path
from typing import Iterable

from pypdf import PdfReader

from .ast_extractor import extract_relations_ast, language_for_suffix
from .config import settings
from .database import transaction, utc_now

TEXT_EXTENSIONS = {
    '.md', '.txt', '.json', '.jsonl', '.yaml', '.yml', '.toml', '.ini', '.env.example',
    '.py', '.js', '.jsx', '.ts', '.tsx', '.css', '.scss', '.html', '.sql', '.sh', '.ps1',
    '.go', '.rs', '.java', '.c', '.h', '.cpp', '.hpp', '.cs', '.php', '.rb', '.swift', '.kt',
}
DOCUMENT_EXTENSIONS = {'.pdf', '.docx', '.xlsx', '.pptx'}
SKIP_DIRS = {
    '.git', '.next', 'node_modules', 'dist', 'build', 'coverage', '.venv', 'venv',
    '__pycache__', '.cache', 'models', 'downloads', 'artifacts', 'generated'
}


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _read_text(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == '.pdf':
        reader = PdfReader(str(path))
        return '\n\n'.join(page.extract_text() or '' for page in reader.pages), 'pdf-text'
    if suffix in {'.docx', '.xlsx', '.pptx'}:
        return _read_office(path, suffix), suffix.lstrip('.')
    data = path.read_bytes()
    if b'\x00' in data[:4096]:
        return '', 'binary'
    for encoding in ('utf-8', 'utf-8-sig', 'latin-1'):
        try:
            return data.decode(encoding), 'text'
        except UnicodeDecodeError:
            continue
    return '', 'unknown'


def _read_office(path: Path, suffix: str) -> str:
    """Extract plain text from .docx / .xlsx / .pptx (all ZIP-based OOXML).

    Raises on a corrupt/non-OOXML file so the indexer records it in the per-file
    errors list instead of silently producing empty content.
    """
    if suffix == '.docx':
        from docx import Document
        doc = Document(str(path))
        parts: list[str] = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text for c in row.cells if c.text]
                if cells:
                    parts.append('\t'.join(cells))
        return '\n'.join(parts)
    if suffix == '.pptx':
        from pptx import Presentation
        deck = Presentation(str(path))
        lines: list[str] = []
        for index, slide in enumerate(deck.slides, 1):
            lines.append(f'# Slide {index}')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text.strip())
        return '\n'.join(lines)
    # suffix == '.xlsx'
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                chunks.append('\t'.join(cells))
    wb.close()
    return '\n'.join(chunks)


def iter_indexable_files(root: Path) -> Iterable[Path]:
    count = 0
    for path in root.rglob('*'):
        if count >= settings.max_index_files:
            break
        if not path.is_file():
            continue
        if any(part in SKIP_DIRS for part in path.parts):
            continue
        suffix = path.suffix.lower()
        if suffix not in TEXT_EXTENSIONS and suffix not in DOCUMENT_EXTENSIONS:
            continue
        if path.stat().st_size > settings.max_file_mb * 1024 * 1024:
            continue
        count += 1
        yield path


def _extract_relations(project_id: int, relative: str, content: str) -> list[tuple]:
    """Build relations table rows for a file via tree-sitter AST extraction.

    Returns 6-tuples (project_id, source_node, relation, target_node, evidence,
    confidence) ready for INSERT OR IGNORE. Files in languages without a
    tree-sitter grammar (Go, Rust, Markdown, JSON, ...) produce no rows — we no
    longer pretend to extract relations from them with fragile regex.
    """
    suffix = Path(relative).suffix.lower()
    language = language_for_suffix(suffix)
    if language is None:
        return []
    ast_rows = extract_relations_ast(relative, content, language)
    return [
        (project_id, relative, relation, target_node, evidence, confidence)
        for relation, target_node, evidence, confidence in ast_rows
    ]


def index_project(project_id: int, root: Path, progress_cb=None) -> dict:
    """Index a project's files into the files/files_fts/relations tables.

    Sets projects.status to 'indexing' on entry, 'ready' on success, 'failed'
    on exception. progress_cb(phase, total, processed) is invoked best-effort
    (never raises) so callers can render live progress; not relied on for
    correctness. Returns the same counts dict as before.

    Transaction layout is deliberate: the 'indexing' state transition and the
    running progress counts are committed in their own short transactions so
    they are VISIBLE to other connections (SQLite WAL readers — the
    /index-status endpoint and the 409 double-index guard) while indexing
    runs. The file loop commits in batches (PROGRESS_COMMIT_EVERY) so live
    progress advances without committing per-file. Partial progress is safe:
    indexing is idempotent (SHA-256 dedup), so a crash mid-run preserves what
    was indexed and the next run picks up the rest.
    """
    PROGRESS_COMMIT_EVERY = 25

    def _report(phase: str, total: int, processed: int) -> None:
        if progress_cb is None:
            return
        try:
            progress_cb(phase, total, processed)
        except Exception:  # progress reporting must never break indexing
            pass

    def _set_progress(project_id: int, phase: str, total: int, processed: int) -> None:
        """Commit a short transaction updating the index_progress row (and,
        when phase is 'indexing' on the first call, flip projects.status).

        Committed immediately so WAL readers observe the current state."""
        ts = utc_now()
        with transaction() as conn:
            conn.execute(
                "INSERT INTO index_progress(project_id,phase,total_files,processed_files,started_at,updated_at,error) "
                "VALUES(?,?,?,?,?,?,?) "
                "ON CONFLICT(project_id) DO UPDATE SET phase=excluded.phase,total_files=excluded.total_files,"
                "processed_files=excluded.processed_files,updated_at=excluded.updated_at,error=NULL",
                (project_id, phase, total, processed, ts, ts, None),
            )

    indexed = 0
    unchanged = 0
    errors: list[str] = []
    seen: set[str] = set()
    now = utc_now()

    try:
        if not root.exists() or not root.is_dir():
            raise FileNotFoundError(f'Project root does not exist: {root}')

        # Transition to 'indexing' in its OWN committed transaction so the
        # state machine is observable to readers before the long loop starts.
        with transaction() as conn:
            conn.execute("UPDATE projects SET status='indexing' WHERE id=?", (project_id,))
            conn.execute(
                "INSERT INTO index_progress(project_id,phase,total_files,processed_files,started_at,updated_at,error) "
                "VALUES(?,?,?,?,?,?,?) "
                "ON CONFLICT(project_id) DO UPDATE SET phase=excluded.phase,total_files=excluded.total_files,"
                "processed_files=excluded.processed_files,started_at=excluded.started_at,updated_at=excluded.updated_at,error=NULL",
                (project_id, 'scanning', 0, 0, now, now, None),
            )

        files_to_index = list(iter_indexable_files(root))
        total = len(files_to_index)
        _set_progress(project_id, 'indexing', total, 0)
        _report('indexing', total, 0)

        since_commit = 0
        for path in files_to_index:
            relative = path.relative_to(root).as_posix()
            seen.add(relative)
            try:
                stat = path.stat()
                raw = path.read_bytes()
                digest = _sha256(raw)
                with transaction() as conn:
                    existing = conn.execute(
                        'SELECT id, sha256 FROM files WHERE project_id=? AND relative_path=?',
                        (project_id, relative)
                    ).fetchone()
                    if existing and existing['sha256'] == digest:
                        unchanged += 1
                    else:
                        content, parser = _read_text(path)
                        if content.strip():
                            if existing:
                                file_id = existing['id']
                                conn.execute(
                                    '''UPDATE files SET absolute_path=?, extension=?, size_bytes=?, modified_ns=?,
                                       sha256=?, content=?, indexed_at=?, parser=?, parse_status=? WHERE id=?''',
                                    (str(path), path.suffix.lower(), stat.st_size, stat.st_mtime_ns,
                                     digest, content, now, parser, 'indexed', file_id)
                                )
                                conn.execute('DELETE FROM files_fts WHERE file_id=?', (file_id,))
                            else:
                                cur = conn.execute(
                                    '''INSERT INTO files(project_id, relative_path, absolute_path, extension, size_bytes,
                                       modified_ns, sha256, content, indexed_at, parser, parse_status) VALUES(?,?,?,?,?,?,?,?,?,?,?)''',
                                    (project_id, relative, str(path), path.suffix.lower(), stat.st_size,
                                     stat.st_mtime_ns, digest, content, now, parser, 'indexed')
                                )
                                file_id = cur.lastrowid
                            conn.execute(
                                'INSERT INTO files_fts(content, relative_path, project_id, file_id) VALUES(?,?,?,?)',
                                (content, relative, project_id, file_id)
                            )
                            conn.execute('DELETE FROM relations WHERE project_id=? AND source_node=?', (project_id, relative))
                            conn.executemany(
                                '''INSERT OR IGNORE INTO relations(project_id,source_node,relation,target_node,evidence,confidence)
                                   VALUES(?,?,?,?,?,?)''',
                                _extract_relations(project_id, relative, content)
                            )
                            indexed += 1
            except Exception as exc:  # keep indexing other files
                errors.append(f'{relative}: {exc}')
            since_commit += 1
            if since_commit >= PROGRESS_COMMIT_EVERY:
                _set_progress(project_id, 'indexing', total, indexed + unchanged)
                _report('indexing', total, indexed + unchanged)
                since_commit = 0

        # Final accounting + state flip in one committed transaction.
        with transaction() as conn:
            existing_paths = [row['relative_path'] for row in conn.execute(
                'SELECT relative_path FROM files WHERE project_id=?', (project_id,)
            )]
            removed = [path for path in existing_paths if path not in seen]
            for relative in removed:
                file_row = conn.execute(
                    'SELECT id FROM files WHERE project_id=? AND relative_path=?', (project_id, relative)
                ).fetchone()
                if file_row:
                    conn.execute('DELETE FROM files_fts WHERE file_id=?', (file_row['id'],))
                conn.execute('DELETE FROM files WHERE project_id=? AND relative_path=?', (project_id, relative))
                conn.execute('DELETE FROM relations WHERE project_id=? AND source_node=?', (project_id, relative))

            conn.execute('UPDATE projects SET indexed_at=?, status=? WHERE id=?', (now, 'ready', project_id))
            conn.execute(
                "UPDATE index_progress SET phase='done',processed_files=?,updated_at=? WHERE project_id=?",
                (indexed + unchanged, utc_now(), project_id),
            )
            conn.execute(
                'INSERT INTO audit_log(project_id,action,detail,created_at) VALUES(?,?,?,?)',
                (project_id, 'project.indexed', json.dumps({'indexed': indexed, 'unchanged': unchanged, 'errors': errors}), now)
            )
            _report('done', total, indexed + unchanged)

        return {'indexed': indexed, 'unchanged': unchanged, 'removed': len(removed), 'errors': errors}
    except Exception as exc:
        fail_now = utc_now()
        with transaction() as conn:
            conn.execute("UPDATE projects SET status='failed' WHERE id=?", (project_id,))
            conn.execute(
                "INSERT INTO index_progress(project_id,phase,total_files,processed_files,started_at,updated_at,error) "
                "VALUES(?,?,?,?,?,?,?) "
                "ON CONFLICT(project_id) DO UPDATE SET phase='failed',error=excluded.error,updated_at=excluded.updated_at",
                (project_id, 'failed', 0, 0, fail_now, fail_now, str(exc)),
            )
            conn.execute(
                'INSERT INTO audit_log(project_id,action,detail,created_at) VALUES(?,?,?,?)',
                (project_id, 'project.index_failed', str(exc)[:300], fail_now)
            )
        raise
